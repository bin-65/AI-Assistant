import streamlit as st
import google.generativeai as genai
from PIL import Image
import pypdf
import docx
import pptx
from docx import Document
from pptx import Presentation
from fpdf import FPDF
import io

# Page configuration
st.set_page_config(
    page_title="AI Multi-Tool Studio", 
    page_icon="⚡", 
    layout="wide",
    initial_sidebar_state="expanded"
)

# ---------------- COLOR THEME SELECTION FILTER ----------------
st.sidebar.markdown("### 🎨 Dashboard Color Filter")
selected_theme = st.sidebar.selectbox(
    "Choose Color Theme",
    ["⚡ Electric Blue", "💜 Neon Purple", "💚 Emerald Green", "🔥 Crimson Red"]
)

# Theme Palette Definitions
themes = {
    "⚡ Electric Blue": {"primary": "#2563eb", "primary_hover": "#1d4ed8", "accent": "#38bdf8", "gradient": "linear-gradient(135deg, #1e2640 0%, #0f172a 100%)"},
    "💜 Neon Purple": {"primary": "#8b5cf6", "primary_hover": "#7c3aed", "accent": "#c084fc", "gradient": "linear-gradient(135deg, #2e1065 0%, #0f172a 100%)"},
    "💚 Emerald Green": {"primary": "#059669", "primary_hover": "#047857", "accent": "#34d399", "gradient": "linear-gradient(135deg, #064e3b 0%, #0f172a 100%)"},
    "🔥 Crimson Red": {"primary": "#dc2626", "primary_hover": "#b91c1c", "accent": "#f87171", "gradient": "linear-gradient(135deg, #450a0a 0%, #0f172a 100%)"}
}

active_theme = themes[selected_theme]

# Custom CSS with Dynamic Theme Injection
st.markdown(f"""
    <style>
    .stApp {{
        background-color: #0b0f17;
        color: #e2e8f0;
    }}
    
    /* Header Card */
    .header-box {{
        background: {active_theme['gradient']};
        padding: 28px;
        border-radius: 16px;
        border: 1px solid #334155;
        margin-bottom: 25px;
        box-shadow: 0 10px 25px -5px rgba(0, 0, 0, 0.4);
    }}
    .main-title {{
        color: #ffffff;
        font-size: 32px;
        font-weight: 800;
        margin: 0;
        letter-spacing: -0.5px;
    }}
    .sub-title {{
        color: #94a3b8;
        font-size: 15px;
        margin-top: 6px;
    }}

    /* Dynamic Buttons */
    .stButton>button {{
        background: {active_theme['primary']} !important;
        color: white !important;
        border: none !important;
        border-radius: 10px !important;
        font-weight: 600 !important;
        padding: 10px 22px !important;
        transition: all 0.3s ease !important;
        box-shadow: 0 4px 14px rgba(0, 0, 0, 0.3) !important;
    }}
    .stButton>button:hover {{
        background: {active_theme['primary_hover']} !important;
        transform: translateY(-1px);
    }}

    /* Dynamic Output Color Labels */
    .copy-header {{
        font-size: 14px;
        font-weight: 700;
        color: {active_theme['accent']};
        display: flex;
        align-items: center;
        gap: 8px;
        margin-bottom: 8px;
    }}
    .download-header {{
        font-size: 14px;
        font-weight: 700;
        color: {active_theme['accent']};
        margin-top: 20px;
        margin-bottom: 10px;
    }}

    /* Tab Active Color */
    .stTabs [aria-selected="true"] {{
        background-color: {active_theme['primary']} !important;
        color: white !important;
    }}
    .stTabs [data-baseweb="tab"] {{
        border-radius: 8px;
        padding: 8px 18px;
        background-color: #1e293b;
        color: #94a3b8;
    }}
    </style>
""", unsafe_allow_html=True)

# Main Banner
st.markdown(f"""
    <div class="header-box">
        <div class="main-title">⚡ AI Multi-Tool Studio</div>
        <div class="sub-title">Enterprise grade multi-purpose AI assistant | Active Theme: <b>{selected_theme}</b></div>
    </div>
""", unsafe_allow_html=True)

# Quick Metric Dashboard Stats
m1, m2, m3, m4 = st.columns(4)
m1.metric("License Access", "Unlimited ♾️")
m2.metric("Processing Speed", "Ultra Fast ⚡")
m3.metric("Supported Formats", "PDF, Word, PPT")
m4.metric("Security Level", "Encrypted 🔒")

st.markdown("---")

# Helper function to generate PDF bytes safely
def generate_pdf_bytes(title, text_content):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", 'B', 14)
    pdf.cell(0, 10, title.encode('latin-1', 'replace').decode('latin-1'), ln=True)
    pdf.set_font("Arial", size=10)
    pdf.ln(5)
    clean_text = text_content.encode('latin-1', 'replace').decode('latin-1')
    pdf.multi_cell(0, 7, clean_text)
    return pdf.output()

# Helper function to generate Word bytes safely
def generate_docx_bytes(title, text_content):
    doc = Document()
    doc.add_heading(title, 0)
    for paragraph in text_content.split("\n\n"):
        doc.add_paragraph(paragraph)
    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()

# Check API Key from Secrets
api_key = st.secrets.get("GEMINI_API_KEY")

if not api_key:
    st.error("🚨 API Key not found! Please configure `GEMINI_API_KEY` in Streamlit Secrets.")
else:
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel("gemini-1.5-flash")

    # Five Active Tabs
    tab1, tab2, tab3, tab4, tab5 = st.tabs([
        "✍️ Content Creator", 
        "🌐 Translator", 
        "⚡ Smart AI Workspace",
        "📚 Academic Writer",
        "📑 Advanced Doc Hub"
    ])

    # ---------------- TAB 1: CONTENT ASSISTANT ----------------
    with tab1:
        st.subheader("✍️ Social Media Content Generator")
        with st.form("content_form"):
            col1, col2 = st.columns(2)
            with col1:
                platform = st.selectbox("Target Platform", ["LinkedIn", "Instagram", "Twitter / X", "Facebook"])
                content_type = st.selectbox("Content Style", ["Informational Post", "Promotional / Ad", "Storytelling"])
            with col2:
                tone = st.selectbox("Tone & Persona", ["Professional", "Casual & Friendly", "Persuasive", "Inspirational"])
                target_audience = st.text_input("Target Audience", placeholder="e.g., Tech Founders, Students")
            
            topic = st.text_area("Core Brief / Topic", placeholder="Describe key talking points...")
            submit_btn = st.form_submit_button("🚀 Generate Content")

        if submit_btn:
            if not topic or not target_audience:
                st.warning("⚠️ Please complete all required fields.")
            else:
                try:
                    prompt = f"Platform: {platform}\nType: {content_type}\nTone: {tone}\nAudience: {target_audience}\nTopic: {topic}"
                    with st.spinner("Generating content..."):
                        response = model.generate_content(prompt)
                    res_text = response.text
                    
                    st.success("🎉 Generated Successfully!")
                    st.markdown("---")
                    
                    st.markdown('<div class="copy-header">📋 One-Click Copy Document:</div>', unsafe_allow_html=True)
                    st.code(res_text, language="markdown")
                    
                    st.markdown('<div class="download-header">📥 Export Options:</div>', unsafe_allow_html=True)
                    col_d1, col_d2 = st.columns(2)
                    with col_d1:
                        st.download_button("📥 Download PDF (.pdf)", data=generate_pdf_bytes(f"{platform} Post", res_text), file_name="Social_Media_Post.pdf", mime="application/pdf", use_container_width=True)
                    with col_d2:
                        st.download_button("📄 Download MS Word (.docx)", data=generate_docx_bytes(f"{platform} Post", res_text), file_name="Social_Media_Post.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document", use_container_width=True)
                except Exception as e:
                    st.error(f"Execution Error: {e}")

    # ---------------- TAB 2: TRANSLATOR ----------------
    with tab2:
        st.subheader("🌐 Global Multi-Language Translator")
        languages_50 = ["English", "Urdu", "Arabic", "Hindi", "Pashto", "Punjabi", "Sindhi", "Spanish", "French", "German", "Chinese"]
        target_language = st.selectbox("Select Target Language", languages_50)
        input_text = st.text_area("Source Text", placeholder="Paste source text here...", height=150)
        translate_btn = st.button("🌐 Translate Content")

        if translate_btn:
            if not input_text.strip():
                st.warning("⚠️ Please provide source text.")
            else:
                try:
                    translation_prompt = f"Automatically detect source language and translate accurately to {target_language}:\n\n{input_text}"
                    with st.spinner("Translating..."):
                        response = model.generate_content(translation_prompt)
                    res_text = response.text
                    
                    st.success("🎉 Translation Complete!")
                    st.markdown("---")
                    
                    st.markdown('<div class="copy-header">📋 One-Click Copy Translation:</div>', unsafe_allow_html=True)
                    st.code(res_text, language="markdown")
                    
                    st.markdown('<div class="download-header">📥 Export Translation:</div>', unsafe_allow_html=True)
                    col_d1, col_d2 = st.columns(2)
                    with col_d1:
                        st.download_button("📥 Download PDF (.pdf)", data=generate_pdf_bytes(f"Translation ({target_language})", res_text), file_name="Translation.pdf", mime="application/pdf", use_container_width=True)
                    with col_d2:
                        st.download_button("📄 Download MS Word (.docx)", data=generate_docx_bytes(f"Translation ({target_language})", res_text), file_name="Translation.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document", use_container_width=True)
                except Exception as e:
                    st.error(f"Execution Error: {e}")

    # ---------------- TAB 3: SMART AI WORKSPACE ----------------
    with tab3:
        st.subheader("⚡ Smart File & Image Workspace")
        user_prompt = st.text_area("Analysis Prompt", placeholder="e.g., Summarize or extract information...", height=100)
        uploaded_files = st.file_uploader(
            "Upload Documents / Images", 
            type=["png", "jpg", "jpeg", "pdf", "docx", "pptx", "txt"], 
            accept_multiple_files=True
        )
        process_btn = st.button("⚡ Process Query")

        if process_btn:
            if not user_prompt.strip() and not uploaded_files:
                st.warning("⚠️ Enter instructions or upload files.")
            else:
                try:
                    with st.spinner("Analyzing data..."):
                        contents = []
                        extracted_text_from_docs = ""

                        if uploaded_files:
                            for file in uploaded_files:
                                file_type = file.type or ""
                                filename = file.name.lower()

                                if "image" in file_type or filename.endswith((".png", ".jpg", ".jpeg")):
                                    contents.append(Image.open(file))
                                elif filename.endswith(".pdf"):
                                    pdf_reader = pypdf.PdfReader(file)
                                    pdf_text = "\n".join([page.extract_text() or "" for page in pdf_reader.pages])
                                    extracted_text_from_docs += f"\n--- [PDF: {file.name}] ---\n{pdf_text}\n"
                                elif filename.endswith(".docx"):
                                    doc_file = docx.Document(file)
                                    docx_text = "\n".join([p.text for p in doc_file.paragraphs])
                                    extracted_text_from_docs += f"\n--- [Word: {file.name}] ---\n{docx_text}\n"
                                elif filename.endswith(".pptx"):
                                    prs_file = pptx.Presentation(file)
                                    pptx_text = ""
                                    for slide in prs_file.slides:
                                        for shape in slide.shapes:
                                            if hasattr(shape, "text"): pptx_text += shape.text + "\n"
                                    extracted_text_from_docs += f"\n--- [PPT: {file.name}] ---\n{pptx_text}\n"
                                elif filename.endswith(".txt"):
                                    extracted_text_from_docs += f"\n--- [TXT: {file.name}] ---\n{file.read().decode('utf-8', errors='ignore')}\n"

                        final_instruction = ""
                        if extracted_text_from_docs:
                            final_instruction += f"=== EXTRACTED DOCUMENTS CONTENT ===\n{extracted_text_from_docs}\n"
                        if user_prompt.strip():
                            final_instruction += f"=== USER INSTRUCTION ===\n{user_prompt}"

                        if final_instruction: contents.append(final_instruction)
                        response = model.generate_content(contents)
                        res_text = response.text

                    st.success("🎉 Analysis Complete!")
                    st.markdown("---")
                    
                    st.markdown('<div class="copy-header">📋 One-Click Copy Result:</div>', unsafe_allow_html=True)
                    st.code(res_text, language="markdown")
                    
                    st.markdown('<div class="download-header">📥 Export Options:</div>', unsafe_allow_html=True)
                    col_d1, col_d2 = st.columns(2)
                    with col_d1:
                        st.download_button("📥 Download PDF (.pdf)", data=generate_pdf_bytes("Workspace Output", res_text), file_name="Smart_Workspace_Output.pdf", mime="application/pdf", use_container_width=True)
                    with col_d2:
                        st.download_button("📄 Download MS Word (.docx)", data=generate_docx_bytes("Workspace Output", res_text), file_name="Smart_Workspace_Output.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document", use_container_width=True)
                except Exception as e:
                    st.error(f"Execution Error: {e}")

    # ---------------- TAB 4: ASSIGNMENT WRITER ----------------
    with tab4:
        st.subheader("📚 Academic & Assignment Writer")
        with st.form("assignment_form"):
            subject_topic = st.text_input("Topic / Subject Header", placeholder="e.g., Deep Learning & Neural Networks")
            academic_level = st.selectbox("Academic Target Level", ["School Level", "High School / College", "Undergraduate", "Postgraduate / PhD"])
            assign_submit = st.form_submit_button("✨ Generate Assignment")

        if assign_submit:
            if not subject_topic.strip():
                st.warning("⚠️ Topic is required.")
            else:
                try:
                    prompt = f"Write a structured, highly detailed academic paper on '{subject_topic}' for {academic_level} level."
                    with st.spinner("Drafting academic content..."):
                        response = model.generate_content(prompt)
                    res_text = response.text
                    
                    st.success("🎉 Assignment Prepared!")
                    st.markdown("---")
                    
                    st.markdown('<div class="copy-header">📋 One-Click Copy Assignment:</div>', unsafe_allow_html=True)
                    st.code(res_text, language="markdown")
                    
                    st.markdown('<div class="download-header">📥 Export Assignment:</div>', unsafe_allow_html=True)
                    col_d1, col_d2 = st.columns(2)
                    with col_d1:
                        st.download_button("📥 Download PDF (.pdf)", data=generate_pdf_bytes(f"Assignment: {subject_topic}", res_text), file_name=f"{subject_topic}_Assignment.pdf", mime="application/pdf", use_container_width=True)
                    with col_d2:
                        st.download_button("📄 Download MS Word (.docx)", data=generate_docx_bytes(f"Assignment: {subject_topic}", res_text), file_name=f"{subject_topic}_Assignment.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document", use_container_width=True)
                except Exception as e:
                    st.error(f"Execution Error: {e}")

    # ---------------- TAB 5: ADVANCED DOCUMENT HUB ----------------
    with tab5:
        st.subheader("📑 Advanced Document & MCQ Hub")
        doc_sub_tab1, doc_sub_tab2 = st.tabs(["📝 Create Document / Manual", "📤 Upload & Extract MCQs"])

        with doc_sub_tab1:
            doc_topic = st.text_input("Document Topic", placeholder="e.g., Fundamentals of Cybersecurity")
            export_format = st.selectbox("Format Structure", ["MS Word (.docx)", "PowerPoint Presentation (.pptx)", "PDF Document (.pdf)"])
            doc_length = st.selectbox("Scope", ["Detailed Notes (~500 words)", "Full Chapter (~1500 words)", "Full Manual (~3000+ words)"])
            custom_prompt = st.text_area("Specific Outline (Optional)")

            create_doc_btn = st.button("✨ Build Document")

            if create_doc_btn:
                if not doc_topic.strip():
                    st.warning("⚠️ Topic required.")
                else:
                    try:
                        with st.spinner("Drafting document..."):
                            gen_prompt = f"Create a structured document on '{doc_topic}' with format target {export_format} and length {doc_length}."
                            response = model.generate_content(gen_prompt)
                            generated_text = response.text

                        st.success("🎉 Document Created!")
                        st.markdown("---")
                        
                        st.markdown('<div class="copy-header">📋 One-Click Copy Document:</div>', unsafe_allow_html=True)
                        st.code(generated_text, language="markdown")

                        st.markdown('<div class="download-header">📥 Export Document:</div>', unsafe_allow_html=True)
                        col_d1, col_d2 = st.columns(2)
                        with col_d1:
                            st.download_button("📥 Download PDF (.pdf)", data=generate_pdf_bytes(doc_topic, generated_text), file_name=f"{doc_topic}.pdf", mime="application/pdf", use_container_width=True)
                        with col_d2:
                            if export_format == "PowerPoint Presentation (.pptx)":
                                prs = Presentation()
                                slides_content = generated_text.split("\n\n")
                                for slide_text in slides_content[:10]:
                                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                                    slide.shapes.title.text = doc_topic
                                    slide.placeholders[1].text = slide_text
                                bio = io.BytesIO()
                                prs.save(bio)
                                st.download_button("📊 Download PPTX (.pptx)", data=bio.getvalue(), file_name=f"{doc_topic}.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)
                            else:
                                st.download_button("📄 Download MS Word (.docx)", data=generate_docx_bytes(doc_topic, generated_text), file_name=f"{doc_topic}.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document", use_container_width=True)

                    except Exception as e:
                        st.error(f"Execution Error: {e}")

        with doc_sub_tab2:
            uploaded_docs = st.file_uploader("Upload Course Material (PDF, Word, PPT, TXT)", type=["pdf", "docx", "pptx", "txt"], accept_multiple_files=True)
            mcq_count = st.selectbox("Question Volume", ["10 MCQs", "20 MCQs", "30 MCQs", "50 MCQs", "100 MCQs"])
            task_type = st.selectbox("Extraction Goal", ["Generate MCQs with Answer Key", "Summarize Key Concepts", "Extract Formulas"])

            process_upload_btn = st.button("🚀 Process & Extract")

            if process_upload_btn:
                if not uploaded_docs:
                    st.warning("⚠️ Upload files first.")
                else:
                    try:
                        with st.spinner("Extracting..."):
                            extracted_full_text = ""
                            for file in uploaded_docs:
                                extracted_full_text += f"\n--- FILE: {file.name} ---\n"
                                if file.name.endswith(".pdf"):
                                    pdf_reader = pypdf.PdfReader(file)
                                    for page in pdf_reader.pages: extracted_full_text += (page.extract_text() or "") + "\n"
                                elif file.name.endswith(".docx"):
                                    doc = docx.Document(file)
                                    for p in doc.paragraphs: extracted_full_text += p.text + "\n"
                                elif file.name.endswith(".pptx"):
                                    prs = pptx.Presentation(file)
                                    for slide in prs.slides:
                                        for shape in slide.shapes:
                                            if hasattr(shape, "text"): extracted_full_text += shape.text + "\n"
                                elif file.name.endswith(".txt"):
                                    extracted_full_text += file.read().decode("utf-8", errors="ignore")

                            mcq_prompt = f"Task: {task_type}\nQuantity: {mcq_count}\nSource Text:\n{extracted_full_text}"
                            response = model.generate_content(mcq_prompt)
                            output_text = response.text

                        st.success("🎉 Extraction Complete!")
                        st.markdown("---")
                        
                        st.markdown('<div class="copy-header">📋 One-Click Copy Result:</div>', unsafe_allow_html=True)
                        st.code(output_text, language="markdown")

                        st.markdown('<div class="download-header">📥 Export Question Bank:</div>', unsafe_allow_html=True)
                        col_mcq1, col_mcq2 = st.columns(2)
                        with col_mcq1:
                            st.download_button("📥 Download PDF (.pdf)", data=generate_pdf_bytes("Extracted MCQs", output_text), file_name="Extracted_MCQs.pdf", mime="application/pdf", use_container_width=True)
                        with col_mcq2:
                            st.download_button("📄 Download MS Word (.docx)", data=generate_docx_bytes("Extracted MCQs", output_text), file_name="Extracted_MCQs.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document", use_container_width=True)

                    except Exception as e:
                        st.error(f"Execution Error: {e}")
